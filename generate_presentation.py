"""Generate MediLink RAG Midterm Presentation (PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# ── Colour palette ──────────────────────────────────────────────
DARK_BG    = RGBColor(0x1B, 0x1F, 0x3B)   # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2    = RGBColor(0x48, 0xCA, 0xE4)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xC0, 0xC0, 0xC0)
MED_GRAY   = RGBColor(0x80, 0x80, 0x80)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)   # success green
RED_SOFT   = RGBColor(0xEF, 0x47, 0x6F)   # soft red
ORANGE     = RGBColor(0xFF, 0xD1, 0x66)   # amber
TRANSPARENT = RGBColor(0x22, 0x26, 0x47)  # slightly lighter navy

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PLOTS = Path("results/plots")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────

def _dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def _add_text(slide, left, top, width, height, text, size=18, color=WHITE,
              bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def _add_para(tf, text, size=16, color=WHITE, bold=False, space_before=Pt(4),
              alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p

def _accent_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), top, Inches(2.2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def _section_title(slide, title, subtitle=""):
    _dark_bg(slide)
    _accent_line(slide, Inches(2.1))
    _add_text(slide, Inches(0.8), Inches(2.25), Inches(11), Inches(1),
              title, size=38, bold=True, color=WHITE)
    if subtitle:
        _add_text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.6),
                  subtitle, size=20, color=LIGHT_GRAY)

def _slide_header(slide, title, number=""):
    _dark_bg(slide)
    label = f"{number}  {title}" if number else title
    _add_text(slide, Inches(0.6), Inches(0.25), Inches(11), Inches(0.7),
              label, size=28, bold=True, color=WHITE)
    # thin accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.95), Inches(12), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def _bullet_slide(slide, items, start_top=Inches(1.35), left=Inches(0.9),
                  size=17, color=WHITE):
    tf = _add_text(slide, left, start_top, Inches(11.5), Inches(5.5),
                   "", size=size, color=color)
    tf.paragraphs[0].text = ""
    for i, item in enumerate(items):
        if i == 0:
            tf.paragraphs[0].text = item
            tf.paragraphs[0].font.size = Pt(size)
            tf.paragraphs[0].font.color.rgb = color
            tf.paragraphs[0].font.name = "Calibri"
            tf.paragraphs[0].space_before = Pt(6)
        else:
            _add_para(tf, item, size=size, color=color, space_before=Pt(6))
    return tf

def _table(slide, rows, col_widths, top, left=Inches(0.6),
           header_color=ACCENT, row_height=Inches(0.38)):
    tbl_shape = slide.shapes.add_table(len(rows), len(col_widths),
                                       left, top,
                                       sum(col_widths), row_height * len(rows))
    tbl = tbl_shape.table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = WHITE
            # shading
            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
            elif ri % 2 == 0:
                fill.fore_color.rgb = RGBColor(0x1E, 0x22, 0x44)
            else:
                fill.fore_color.rgb = TRANSPARENT
    return tbl

def _add_plot(slide, path, left, top, width=None, height=None):
    p = Path(path)
    if p.exists():
        kwargs = {}
        if width:  kwargs["width"] = width
        if height: kwargs["height"] = height
        slide.shapes.add_picture(str(p), left, top, **kwargs)
        return True
    return False


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_dark_bg(sl)
_add_text(sl, Inches(1), Inches(1.6), Inches(11), Inches(1.2),
          "MediLink RAG", size=52, bold=True, color=ACCENT)
_add_text(sl, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
          "Bilingual Medical Retrieval-Augmented Generation System",
          size=26, color=WHITE)
_accent_line(sl, Inches(3.65))
_add_text(sl, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
          "Midterm Presentation  |  March 2026", size=18, color=LIGHT_GRAY)
_add_text(sl, Inches(1), Inches(4.55), Inches(11), Inches(0.5),
          "Hardware: Multi-GPU RTX A6000 (47 GB VRAM each)", size=15, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Problem Statement", "01")
_bullet_slide(sl, [
    "Arabic is underserved by most LLMs — no reliable, grounded medical Q&A exists",
    "General-purpose LLMs hallucinate medical facts, especially in Arabic",
    "Patients and doctors need trustworthy, source-cited answers",
    "Challenge: English medical corpus + Arabic user queries = cross-lingual retrieval",
], start_top=Inches(1.4), size=19)

tf = _add_text(sl, Inches(0.9), Inches(4.3), Inches(11), Inches(1.5), "", size=17)
_add_para(tf, "Our Solution", size=22, bold=True, color=ACCENT)
_add_para(tf, "A RAG pipeline that retrieves from an English medical textbook,", size=17, color=LIGHT_GRAY)
_add_para(tf, "generates grounded answers in Arabic, and validates every response", size=17, color=LIGHT_GRAY)
_add_para(tf, "through multi-layer safety checks before delivery.", size=17, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview (condensed text version)
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "End-to-End Pipeline", "02")

steps = [
    ("1", "Emergency Detection",  "30+ AR/EN keywords → immediate escalation"),
    ("2", "Preprocessing",        "Arabic normalization (alef, ya, diacritics)"),
    ("3", "Query Expansion",      "Medical dictionary → English synonym variants"),
    ("4", "Query Translation",    "Llama-3.1-8B (Groq) Arabic→English for BM25 path"),
    ("5", "Dense Retrieval",      "BGE-M3 → FAISS top-20 per variant"),
    ("6", "BM25 Retrieval",       "English-translated query → BM25-Okapi top-20"),
    ("7", "Hybrid Fusion",        "Intent-aware weighted merge + agreement boost"),
    ("8", "Reranking (optional)", "BGE-Reranker-v2-M3 cross-encoder"),
    ("9", "Context Building",     "Score > 0.25, budget 1500 tokens, top 5-7 chunks"),
    ("10", "LLM Generation",     "Groq → Llama-3.1-8B, temp=0.2"),
    ("11", "Safety Validation",   "Content filter + LLM Judge (grounding < 0.3 → reject)"),
]

left_col_x = Inches(0.6)
right_col_x = Inches(6.6)
y_start = Inches(1.35)
row_h = Inches(0.46)

for i, (num, title, desc) in enumerate(steps):
    col_x = left_col_x if i < 6 else right_col_x
    row_i = i if i < 6 else i - 6
    y = y_start + row_i * row_h

    # number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, col_x, y, Inches(0.38), Inches(0.38))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    _add_text(sl, col_x + Inches(0.5), y - Pt(2), Inches(5.2), Inches(0.24),
              title, size=15, bold=True, color=WHITE)
    _add_text(sl, col_x + Inches(0.5), y + Inches(0.2), Inches(5.2), Inches(0.24),
              desc, size=12, color=LIGHT_GRAY)

_add_text(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
          "Total latency: ~4-6 s  (generation ~2 s + judge ~2-3 s)",
          size=14, color=ORANGE, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Indexing Pipeline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Indexing Pipeline (Offline)", "03")

idx_rows = [
    ["Step", "Method", "Detail"],
    ["PDF Loading", "PyPDF2", "Page-by-page text extraction"],
    ["Preprocessing", "Custom", "Arabic normalization (alef, ya, ta-marbuta, diacritics)"],
    ["Tokenization", "Custom Stemmer", "~50 Arabic stopwords, bilingual auto-detection"],
    ["Chunking", "Semantic", "200-word chunks, 40-word overlap"],
    ["Embedding", "BAAI/bge-m3", "1024-dim, multilingual, L2-normalized, batch=64"],
    ["Vector Store", "FAISS IndexFlatIP", "Cosine similarity via inner product"],
    ["BM25 Index", "BM25-Okapi", "Separate AR/EN tokenizers"],
]
_table(sl, idx_rows,
       [Inches(1.8), Inches(2.2), Inches(7.5)],
       top=Inches(1.35))

_add_text(sl, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
          "Result: 759 indexed chunks from Gale Encyclopedia of Medicine",
          size=16, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Retrieval Pipeline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Retrieval Pipeline", "04")

ret_rows = [
    ["Component", "Method", "Key Detail"],
    ["Query Expansion", "Medical dictionary", "~100+ terms → 2-3 English variants"],
    ["Query Translation", "Groq LLM", "Arabic→English for BM25 path"],
    ["Dense Retrieval", "BGE-M3 + FAISS", "Multi-query embedding → top-20 per variant"],
    ["BM25 Retrieval", "BM25-Okapi", "English-translated query → top-20"],
    ["Hybrid Fusion", "Intent-aware", "Dense 0.8 / BM25 0.2 + agreement boost 1.2×"],
    ["Reranking", "BGE-Reranker-v2-M3", "Cross-encoder (disabled by default)"],
    ["Dynamic Selection", "Score + budget", "Threshold > 0.25 → top 5-7 chunks, 1500 tokens"],
]
_table(sl, ret_rows,
       [Inches(2.2), Inches(2.5), Inches(7.0)],
       top=Inches(1.35))

_add_text(sl, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
          "Reranker disabled: degraded Arabic Recall@1 to 0.004 in early experiments",
          size=14, color=RED_SOFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Generation & Safety
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Generation & Safety Pipeline", "05")

# Generation side
_add_text(sl, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4),
          "Generation", size=22, bold=True, color=ACCENT)
_bullet_slide(sl, [
    "Strict prompt: answer ONLY from retrieved context",
    "Max 6 bullet points, role-adapted (patient vs. doctor)",
    "Groq API → Llama-3.1-8B-Instant (temp=0.2)",
    "Post-processing: line deduplication + sentence truncation",
], start_top=Inches(1.85), left=Inches(0.8), size=15)

# Safety side
_add_text(sl, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
          "Multi-Layer Safety", size=22, bold=True, color=ACCENT)

safety_items = [
    ("Emergency Detector", "30+ AR/EN keywords → escalation"),
    ("Content Filter", "PII regex (SSN, credit card, phone, email)"),
    ("LLM Judge", "Grounding + hallucination + confidence scoring"),
    ("Rejection Gate", "Grounding < 0.3 → answer rejected"),
]
y = Inches(1.85)
for title, desc in safety_items:
    _add_text(sl, Inches(6.8), y, Inches(5.5), Inches(0.25),
              title, size=15, bold=True, color=WHITE)
    _add_text(sl, Inches(6.8), y + Inches(0.25), Inches(5.5), Inches(0.25),
              desc, size=13, color=LIGHT_GRAY)
    y += Inches(0.6)

# Bottom: calibration note
_add_text(sl, Inches(0.6), Inches(5.5), Inches(12), Inches(1.0),
          "Confidence Calibration: Logistic Regression on 6 features (grounding, retrieval, "
          "rerank, context_len, answer_len, top_similarity). Fallback heuristic when < 40 samples.",
          size=13, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Models Used
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Models & Technologies", "06")

model_rows = [
    ["Model", "Role", "Dim", "Provider"],
    ["BAAI/bge-m3", "Dense Embedding (AR+EN)", "1024", "HuggingFace (local GPU)"],
    ["BAAI/bge-reranker-v2-m3", "Cross-encoder Reranking", "—", "HuggingFace (local GPU)"],
    ["Llama-3.1-8B-Instant", "Generation + Judge + Translation", "—", "Groq API (cloud)"],
    ["Qwen2.5-32B-Instruct", "Evaluation-only Generation", "—", "HuggingFace (local GPU)"],
]
_table(sl, model_rows,
       [Inches(3.2), Inches(3.8), Inches(1.0), Inches(3.5)],
       top=Inches(1.5))

tech_items = [
    "FAISS (IndexFlatIP)  •  BM25-Okapi (rank_bm25)  •  FastAPI + Uvicorn",
    "PyPDF2  •  Sentence Transformers  •  scikit-learn (calibration)  •  Matplotlib / Seaborn",
]
tf = _add_text(sl, Inches(0.6), Inches(4.8), Inches(12), Inches(1.5),
               "Tech Stack", size=20, bold=True, color=ACCENT)
for t in tech_items:
    _add_para(tf, t, size=15, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Evaluation: Retrieval Metrics
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Retrieval Evaluation (99 Queries, 4 Modes)", "07")

eval_rows = [
    ["Mode", "Recall@1", "Recall@5", "Recall@10", "MRR", "nDCG@10", "Hit@10", "Zero-Recall"],
    ["BM25",           "0.06", "0.20", "0.29", "0.41", "0.27", "0.66", "34"],
    ["Dense",          "0.20", "0.87", "0.97", "0.98", "0.96", "0.99", "1"],
    ["Hybrid",         "0.18", "0.74", "0.89", "0.92", "0.87", "0.98", "2"],
    ["Hybrid+Rerank",  "0.12", "0.50", "0.89", "0.75", "0.74", "0.98", "2"],
]
_table(sl, eval_rows,
       [Inches(1.8), Inches(1.1), Inches(1.1), Inches(1.2), Inches(0.9),
        Inches(1.2), Inches(1.1), Inches(1.4)],
       top=Inches(1.35), row_height=Inches(0.42))

# Caveat box
box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(4.2), Inches(11.8), Inches(2.6))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x33, 0x18, 0x18)
box.line.color.rgb = RED_SOFT
box.line.width = Pt(1.5)

tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "⚠  Ground Truth Caveat"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RED_SOFT
_add_para(tf, "The ground truth was auto-generated by running dense retrieval and assigning its "
          "top-5 results as 'relevant docs'. This is circular evaluation — dense metrics are "
          "inflated because the system recovers its own previous outputs.", size=13, color=LIGHT_GRAY)
_add_para(tf, "These numbers measure retrieval self-consistency, not true retrieval quality. "
          "Proper evaluation requires human-annotated relevance judgments.", size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation Observations
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Key Observations", "08")

observations = [
    ("BM25 is weakest (34/99 zero-recall)",
     "Corpus is English; BM25 cannot cross the Arabic→English language barrier "
     "even with translation."),
    ("Hybrid fusion partially recovers BM25",
     "But dilutes dense signal — Recall@5 drops from 0.87 → 0.74."),
    ("Reranking degrades results",
     "Cross-encoder may not generalize to Arabic queries vs. English chunks."),
    ("Dense appears strongest",
     "Partly an artifact of circular ground truth — true quality needs independent validation."),
    ("Generation grounding tracks retrieval quality",
     "Where retrieval fails (low recall), grounding falls and the judge rejects the answer."),
]

y = Inches(1.4)
for title, desc in observations:
    _add_text(sl, Inches(0.9), y, Inches(11), Inches(0.3),
              title, size=17, bold=True, color=ACCENT2)
    _add_text(sl, Inches(0.9), y + Inches(0.32), Inches(11), Inches(0.35),
              desc, size=14, color=LIGHT_GRAY)
    y += Inches(0.85)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Retrieval Quality Plot
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Retrieval Quality — Visual", "09")
if not _add_plot(sl, PLOTS / "retrieval_quality.png",
                 Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5.5)):
    if not _add_plot(sl, PLOTS / "2_retrieval_quality.png",
                     Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5.5)):
        _add_text(sl, Inches(3), Inches(3.5), Inches(7), Inches(1),
                  "[retrieval_quality.png not found]", size=20, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Category Heatmap
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Recall by Category", "10")
if not _add_plot(sl, PLOTS / "retrieval_category_heatmap.png",
                 Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5.5)):
    _add_text(sl, Inches(3), Inches(3.5), Inches(7), Inches(1),
              "[retrieval_category_heatmap.png not found]", size=20, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — MRR Boxplot
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "MRR Distribution by Mode", "11")
if not _add_plot(sl, PLOTS / "retrieval_mrr_boxplot.png",
                 Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5.5)):
    if not _add_plot(sl, PLOTS / "6_mrr_by_mode.png",
                     Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5.5)):
        _add_text(sl, Inches(3), Inches(3.5), Inches(7), Inches(1),
                  "[mrr_boxplot.png not found]", size=20, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Design Decisions
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Key Design Decisions", "12")

decisions = [
    ("BGE-M3 over monolingual models",
     "Best multilingual embedding for Arabic+English; 1024-dim captures medical semantics"),
    ("Groq API over local LLM",
     "Fast inference (~2 s), saves GPU memory — trade-off: 500K tokens/day free tier"),
    ("LLM Judge over cosine grounding",
     "Can reason about paraphrasing and Arabic synonyms; cosine gave false negatives"),
    ("Dense-primary fusion (0.8 / 0.2)",
     "BM25 cannot cross the language barrier; dense handles it natively"),
    ("Disabled reranker by default",
     "Cross-encoder degraded Arabic Recall@1 to 0.004 in experiments"),
    ("Overlapping chunks (40-word overlap)",
     "Medical concepts span paragraph boundaries; prevents information loss"),
]

y = Inches(1.4)
for title, desc in decisions:
    _add_text(sl, Inches(0.9), y, Inches(11), Inches(0.28),
              title, size=16, bold=True, color=GREEN)
    _add_text(sl, Inches(0.9), y + Inches(0.3), Inches(11), Inches(0.3),
              desc, size=13, color=LIGHT_GRAY)
    y += Inches(0.76)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Configuration
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "System Configuration", "13")

conf_rows = [
    ["Parameter", "Value", "Purpose"],
    ["Embedding", "BAAI/bge-m3 (1024-d)", "Multilingual dense retrieval"],
    ["Reranker", "bge-reranker-v2-m3", "Cross-encoder (disabled)"],
    ["LLM", "Llama-3.1-8B (Groq)", "Generation + Judge + Translation"],
    ["Dense / BM25 weight", "0.80 / 0.20", "Hybrid fusion balance"],
    ["Top-K (dense / BM25)", "20 / 20", "Initial retrieval candidates"],
    ["Top-K final", "10", "After fusion + reranking"],
    ["Similarity threshold", "0.25", "Min score for selection"],
    ["Chunk size / overlap", "200 / 40 words", "Semantic chunking"],
    ["Max context tokens", "1500", "LLM prompt budget"],
    ["Judge reject threshold", "< 0.3 grounding", "Answer rejection gate"],
    ["Rate limit", "60 req/min/IP", "API protection"],
]
_table(sl, conf_rows,
       [Inches(2.8), Inches(3.0), Inches(5.5)],
       top=Inches(1.35), row_height=Inches(0.4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Testing & Evaluation Pipeline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Testing & Evaluation Infrastructure", "14")

_add_text(sl, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4),
          "Unit Tests (7 files)", size=20, bold=True, color=ACCENT)
test_items = [
    "test_retrieval.py — fusion, RRF, reranker, expansion",
    "test_safety.py — emergency detection, content filter, PII",
    "test_indexing.py — embedder shape, vector store, BM25",
    "test_calibration.py — ECE, model training, synthetic data",
    "test_api.py — /health, /ready, /ask endpoints",
    "test_arabic_evaluation.py — dataset structure validation",
    "test_calibration_data.py — logging, labeling, records",
]
_bullet_slide(sl, test_items, start_top=Inches(1.8), left=Inches(0.8), size=14, color=LIGHT_GRAY)

_add_text(sl, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
          "Evaluation Scripts", size=20, bold=True, color=ACCENT)
eval_items = [
    ("evaluate_retrieval.py", "4 modes × 99 queries → CSV + plots"),
    ("evaluate_plots.py", "End-to-end with Qwen2.5-32B (local)"),
    ("Metrics", "Recall@k, MRR, nDCG@k, ECE, grounding rate"),
]
y = Inches(1.8)
for title, desc in eval_items:
    _add_text(sl, Inches(6.8), y, Inches(5.5), Inches(0.25),
              title, size=15, bold=True, color=WHITE)
    _add_text(sl, Inches(6.8), y + Inches(0.25), Inches(5.5), Inches(0.25),
              desc, size=13, color=LIGHT_GRAY)
    y += Inches(0.6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Next Steps
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Recommendations & Next Steps", "15")

recs = [
    ("Human-annotated ground truth",
     "Replace auto-generated GT with expert relevance judgments for valid evaluation"),
    ("Upgrade Groq API tier",
     "Free tier (500K tokens/day) is insufficient for production or heavy evaluation"),
    ("Consider dense-only mode as default",
     "Hybrid fusion dilutes dense signal without meaningful BM25 benefit for Arabic"),
    ("Fine-tune reranker on Arabic medical data",
     "Or remove it — current cross-encoder degrades Arabic ranking"),
    ("Add Arabic medical corpus",
     "English-only source limits BM25; Arabic literature would reduce translation dependency"),
    ("Implement query caching",
     "Reduce API calls and latency for repeated queries"),
]

y = Inches(1.4)
for i, (title, desc) in enumerate(recs):
    num_shape = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.7), y + Pt(2), Inches(0.35), Inches(0.35))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = ACCENT
    num_shape.line.fill.background()
    ntf = num_shape.text_frame
    ntf.paragraphs[0].text = str(i + 1)
    ntf.paragraphs[0].font.size = Pt(12)
    ntf.paragraphs[0].font.bold = True
    ntf.paragraphs[0].font.color.rgb = DARK_BG
    ntf.paragraphs[0].alignment = PP_ALIGN.CENTER

    _add_text(sl, Inches(1.2), y, Inches(11), Inches(0.28),
              title, size=16, bold=True, color=WHITE)
    _add_text(sl, Inches(1.2), y + Inches(0.3), Inches(11), Inches(0.28),
              desc, size=13, color=LIGHT_GRAY)
    y += Inches(0.72)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Thank You
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_dark_bg(sl)
_add_text(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
          "Thank You", size=48, bold=True, color=ACCENT,
          alignment=PP_ALIGN.CENTER)
_accent_line(sl, Inches(3.55))
_add_text(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
          "Questions?", size=24, color=WHITE, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────
out = Path("MediLink_RAG_Midterm.pptx")
prs.save(str(out))
print(f"Saved → {out}  ({out.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")
