#!/usr/bin/env python3
"""Update MediLink presentation with honest (LLM-annotated) retrieval metrics."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PPTX_PATH = "MediLink_RAG_Midterm.pptx"

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0x66, 0x66, 0x66)
WARN_AMBER = RGBColor(0xE6, 0x9F, 0x00)


def _set_cell(cell, text, font_size=11, bold=False, color=DARK, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_text_in_shape(shape, text, font_size=None, bold=None, color=None, italic=None):
    """Replace all text in shape's text frame."""
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.text = ""
    # Set first paragraph first run
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if italic is not None:
        run.font.italic = italic


def main():
    prs = Presentation(PPTX_PATH)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 8 (index 7): Update retrieval metrics table
    # ══════════════════════════════════════════════════════════════
    slide8 = prs.slides[7]

    # New metrics from honest evaluation
    # Mode: [Recall@1, Recall@5, Recall@10, MRR, nDCG@10, Hit@10, Zero-Recall]
    new_data = {
        "BM25":           ["0.02", "0.11", "0.24", "0.18", "0.15", "0.44", "45"],
        "Dense":          ["0.06", "0.21", "0.37", "0.35", "0.27", "0.62", "31"],
        "Hybrid":         ["0.04", "0.21", "0.34", "0.31", "0.25", "0.58", "34"],
        "Hybrid+Rerank":  ["0.04", "0.23", "0.34", "0.29", "0.25", "0.58", "34"],
    }

    for shape in slide8.shapes:
        if shape.has_table:
            table = shape.table
            mode_order = ["BM25", "Dense", "Hybrid", "Hybrid+Rerank"]
            for ri, mode in enumerate(mode_order, start=1):
                vals = new_data[mode]
                # Column 0 = mode name (keep as-is)
                for ci, val in enumerate(vals):
                    _set_cell(table.rows[ri].cells[ci + 1], val, font_size=11, bold=False)
                # Bold the mode name
                _set_cell(table.rows[ri].cells[0], mode, font_size=11, bold=True,
                          alignment=PP_ALIGN.LEFT)
            print("✓ Updated Slide 8 metrics table")
            break

    # Update the GT caveat box
    for shape in slide8.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs)
            if "Ground Truth Caveat" in full_text:
                tf = shape.text_frame
                # Paragraph 0: title (keep)
                # Paragraph 1: description
                # Paragraph 2: note
                if len(tf.paragraphs) >= 2:
                    p1 = tf.paragraphs[1]
                    for run in p1.runs:
                        run.text = ""
                    if p1.runs:
                        p1.runs[0].text = (
                            "Ground truth independently annotated by Qwen2.5-32B-Instruct "
                            "(LLM-as-judge). For each of 99 queries, a 2-stage pipeline "
                            "(embedding pre-filter → LLM judge) identified truly relevant "
                            "chunks (avg 2.4 per query)."
                        )
                if len(tf.paragraphs) >= 3:
                    p2 = tf.paragraphs[2]
                    for run in p2.runs:
                        run.text = ""
                    if p2.runs:
                        p2.runs[0].text = (
                            "These are honest metrics — no circular bias. "
                            "Dense retrieval is still strongest but realistic numbers "
                            "reveal substantial room for improvement."
                        )
                print("✓ Updated Slide 8 GT caveat")
                break

    # ══════════════════════════════════════════════════════════════
    # SLIDE 9 (index 8): Update Key Observations
    # ══════════════════════════════════════════════════════════════
    slide9 = prs.slides[8]

    # Map old observation text → new text
    updates = [
        # Observation 1 title
        ("BM25 is weakest (34/99 zero-recall)",
         "BM25 is weakest (45/81 zero-recall)"),
        # Observation 1 detail
        ("Corpus is English; BM25 cannot cross the Arabic→English language barrier even with translation.",
         "Corpus is English; BM25 cannot cross the Arabic→English barrier. Only 44% hit rate."),
        # Observation 2 title
        ("Hybrid fusion partially recovers BM25",
         "Hybrid fusion does not help over dense-only"),
        # Observation 2 detail
        ("But dilutes dense signal — Recall@5 drops from 0.87 → 0.74.",
         "Dilutes dense signal — Recall@10 drops from 0.37 → 0.34. BM25 noise hurts."),
        # Observation 3 title (keep same)
        # Observation 3 detail (keep same)
        # Observation 4 title
        ("Dense appears strongest",
         "Dense is strongest (Recall@10 = 0.37, MRR = 0.35)"),
        # Observation 4 detail
        ("Partly an artifact of circular ground truth — true quality needs independent validation.",
         "Confirmed by independent LLM-judge annotation — no circular bias. Still significant room for improvement."),
    ]

    for shape in slide9.shapes:
        if shape.has_text_frame:
            for old_text, new_text in updates:
                full = shape.text_frame.paragraphs[0].text.strip() if shape.text_frame.paragraphs else ""
                if full == old_text:
                    for run in shape.text_frame.paragraphs[0].runs:
                        run.text = ""
                    if shape.text_frame.paragraphs[0].runs:
                        shape.text_frame.paragraphs[0].runs[0].text = new_text
                    break

    print("✓ Updated Slide 9 Key Observations")

    # ══════════════════════════════════════════════════════════════
    # SLIDES 10-12: Replace plot images with new ones
    # ══════════════════════════════════════════════════════════════
    plot_slides = {
        9: "results/plots/retrieval_quality.png",      # Slide 10: Retrieval Quality
        10: "results/plots/retrieval_category_heatmap.png",  # Slide 11: Recall by Category
        11: "results/plots/retrieval_mrr_boxplot.png",  # Slide 12: MRR Distribution
    }

    for slide_idx, plot_path in plot_slides.items():
        slide = prs.slides[slide_idx]
        # Find the image shape and replace it
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                # Get position and size
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                # Remove old image
                sp = shape._element
                sp.getparent().remove(sp)
                # Add new image at same position
                slide.shapes.add_picture(plot_path, left, top, width, height)
                print(f"✓ Replaced plot on Slide {slide_idx + 1}")
                break

    # ══════════════════════════════════════════════════════════════
    # Save
    # ══════════════════════════════════════════════════════════════
    prs.save(PPTX_PATH)
    print(f"\n✓ Saved updated presentation: {PPTX_PATH}")


if __name__ == "__main__":
    main()
