#!/usr/bin/env python3
"""Fix MediLink presentation: update Slide 4, add Generation Eval & Example slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PPTX_PATH = "MediLink_RAG_Midterm.pptx"
OUTPUT_PATH = "MediLink_RAG_Midterm.pptx"  # overwrite

# ── Colours & fonts (match existing slides) ──────────────────────────
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)  # teal-ish for headers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0x66, 0x66, 0x66)
WARN_AMBER = RGBColor(0xE6, 0x9F, 0x00)
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xF8)


def _set_cell(cell, text, font_size=11, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """Helper: set text in a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _header_cell(cell, text):
    _set_cell(cell, text, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set(qn("a:fill"), "")
    # Set background
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": "1A1A2E"})
    solidFill.append(srgb)
    tcPr.append(solidFill)


def _shade_row(row, color_hex="F0F4F8"):
    from pptx.oxml.ns import qn
    for cell in row.cells:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": color_hex})
        solidFill.append(srgb)
        tcPr.append(solidFill)


# ── Slide helpers ─────────────────────────────────────────────────────
def _add_title(slide, number, title_text):
    """Add the standard title bar (number  Title) + underline."""
    txBox = slide.shapes.add_textbox(Emu(548640), Emu(228600), Emu(10058400), Emu(640080))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{number}  {title_text}"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = DARK
    # Underline bar
    slide.shapes.add_shape(
        1, Emu(548640), Emu(868680), Emu(10972800), Emu(25400)
    )


def _add_footnote(slide, text, top=5486400):
    txBox = slide.shapes.add_textbox(Emu(548640), Emu(top), Emu(10972800), Emu(457200))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = LIGHT_GREY


# ═══════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation(PPTX_PATH)

    # ── 1. Fix Slide 4 (index 3): "Custom Stemmer" → "Bilingual Tokenizer" ──
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_table:
            table = shape.table
            # Row 3 (Tokenization) – fix Method column
            _set_cell(table.rows[3].cells[1], "Bilingual Tokenizer", bold=True, font_size=11)
            print("✓ Fixed Slide 4: Tokenization method → Bilingual Tokenizer")
            break

    # ── 2. Add "Generation Evaluation" slide after Slide 11 (MRR Distribution) ──
    # Slide order: 1-12 existing, insert new at index 12 (after MRR, before Key Design Decisions)
    # Actually let's insert after Slide 9 (Key Observations) to keep eval slides together
    # Current order: 8=Retrieval Eval, 9=Key Obs, 10=Viz, 11=Recall Cat, 12=MRR
    # Insert after slide 12 (MRR Distribution, index 11)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # --- Generation Eval slide ---
    gen_slide = prs.slides.add_slide(blank_layout)
    _add_title(gen_slide, "12", "Generation Evaluation (Qwen2.5-32B Judge)")

    # Table: 4 rows data + 1 header = 5 rows × 5 cols
    rows, cols = 5, 5
    tbl_shape = gen_slide.shapes.add_table(
        rows, cols,
        Emu(548640), Emu(1234440),
        Emu(10515600), Emu(1920240)
    )
    table = tbl_shape.table

    headers = ["Mode", "Grounding", "Confidence", "Retrieval Score", "Grounded %"]
    for ci, h in enumerate(headers):
        _header_cell(table.rows[0].cells[ci], h)

    data = [
        ["Dense",          "0.645", "0.539", "0.576", "59.6%"],
        ["Hybrid",         "0.661", "0.569", "0.581", "62.6%"],
        ["Hybrid+Rerank",  "0.712", "0.666", "0.758", "70.7%"],
    ]
    for ri, row_data in enumerate(data, start=1):
        for ci, val in enumerate(row_data):
            _set_cell(table.rows[ri].cells[ci], val, font_size=11,
                      bold=(ci == 0), alignment=PP_ALIGN.CENTER)
        if ri % 2 == 0:
            _shade_row(table.rows[ri])

    # Overall row
    overall = ["Overall (297)", "0.673", "0.591", "0.638", "64.3%"]
    for ci, val in enumerate(overall):
        _set_cell(table.rows[4].cells[ci], val, font_size=11, bold=True,
                  alignment=PP_ALIGN.CENTER)
    _shade_row(table.rows[4], "E8F4FD")

    # Observations below table
    obs_box = gen_slide.shapes.add_textbox(
        Emu(548640), Emu(3400000), Emu(10515600), Emu(2000000)
    )
    tf = obs_box.text_frame
    tf.word_wrap = True

    bullets = [
        "Hybrid+Rerank achieves highest grounding (0.712) and confidence (0.666)",
        "English queries score higher (avg 0.729) than Arabic (avg 0.617) — cross-lingual gap persists",
        "64.3% of answers are grounded — remaining 35.7% rejected by safety gate",
        "Evaluator: Qwen2.5-32B-Instruct (local GPU), 99 queries × 3 modes = 297 evaluations",
    ]
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(13)
        run.font.color.rgb = DARK

    _add_footnote(gen_slide,
                  "Grounding scored by LLM judge (0-1). Grounded = score ≥ 0.5.",
                  top=5700000)

    print("✓ Added Generation Evaluation slide")

    # --- Concrete Example slide ---
    ex_slide = prs.slides.add_slide(blank_layout)
    _add_title(ex_slide, "13", "Concrete Example — Arabic Medical Q&A")

    # Query box
    q_box = ex_slide.shapes.add_textbox(
        Emu(548640), Emu(1100000), Emu(10515600), Emu(320000)
    )
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "User Query"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    q_text_box = ex_slide.shapes.add_textbox(
        Emu(548640), Emu(1400000), Emu(10515600), Emu(400000)
    )
    tf = q_text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "ما هي أعراض مرض السكري؟"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = "(What are the symptoms of diabetes?)"
    run2.font.size = Pt(12)
    run2.font.italic = True
    run2.font.color.rgb = LIGHT_GREY

    # Retrieved Context box
    ctx_label = ex_slide.shapes.add_textbox(
        Emu(548640), Emu(2100000), Emu(5029200), Emu(280000)
    )
    tf = ctx_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Retrieved Context (top chunks)"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    ctx_box = ex_slide.shapes.add_textbox(
        Emu(548640), Emu(2400000), Emu(5029200), Emu(2800000)
    )
    tf = ctx_box.text_frame
    tf.word_wrap = True
    chunks = [
        "Chunk 1 (score 0.82): \"Diabetes mellitus ... symptoms include frequent urination (polyuria), excessive thirst (polydipsia), and unexplained weight loss...\"",
        "Chunk 2 (score 0.78): \"Type 2 diabetes may also present with blurred vision, slow wound healing, and recurrent infections...\"",
        "Chunk 3 (score 0.71): \"Diabetic ketoacidosis symptoms: nausea, vomiting, abdominal pain, fruity breath odor...\"",
    ]
    for i, chunk in enumerate(chunks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(10)
        run.font.color.rgb = DARK

    # Generated Answer box
    ans_label = ex_slide.shapes.add_textbox(
        Emu(6000000), Emu(2100000), Emu(5200000), Emu(280000)
    )
    tf = ans_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Generated Answer"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    ans_box = ex_slide.shapes.add_textbox(
        Emu(6000000), Emu(2400000), Emu(5200000), Emu(2800000)
    )
    tf = ans_box.text_frame
    tf.word_wrap = True
    answer_lines = [
        "أعراض مرض السكري تشمل:",
        "• كثرة التبول (Polyuria)",
        "• العطش الشديد (Polydipsia)",
        "• فقدان الوزن غير المبرر",
        "• عدم وضوح الرؤية",
        "• بطء التئام الجروح",
        "• التهابات متكررة",
    ]
    for i, line in enumerate(answer_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = DARK
        if i == 0:
            run.font.bold = True

    # Safety verdict
    verdict_box = ex_slide.shapes.add_textbox(
        Emu(548640), Emu(5400000), Emu(10515600), Emu(500000)
    )
    tf = verdict_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Safety Verdict:  Grounding = 0.9  |  Confidence = 0.9  |  ✓ Answer Accepted"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x80, 0x00)

    _add_footnote(ex_slide,
                  "Mode: hybrid_rerank  |  LLM: Llama-3.1-8B-Instant (Groq)  |  Latency: ~4 s",
                  top=5900000)

    print("✓ Added Concrete Example slide")

    # ── 3. Move new slides to correct positions ──
    # New slides are appended at the end (indices 17, 18)
    # We want them after slide 12 (MRR Distribution, index 11)
    # Move generation eval slide (currently last-1) to position after index 11
    # Move example slide (currently last) to position after that

    slide_list = prs.slides._sldIdLst
    slides = list(slide_list)

    # The two new slides are the last two
    gen_eval_elem = slides[-2]
    example_elem = slides[-1]

    # Remove them from current position
    slide_list.remove(gen_eval_elem)
    slide_list.remove(example_elem)

    # Insert after slide 12 (index 11) — which is MRR Distribution
    # After removal, the remaining slides are 0..16 (original 17)
    # We want to insert after index 11
    remaining = list(slide_list)
    ref_elem = remaining[12]  # element at index 12 (Key Design Decisions)

    slide_list.insert(slide_list.index(ref_elem), gen_eval_elem)
    slide_list.insert(slide_list.index(ref_elem), example_elem)

    print("✓ Moved new slides to correct position (after MRR Distribution)")

    # ── 4. Renumber titles for slides that shifted ──
    # After insertion: slides 0-11 unchanged, 12=Gen Eval, 13=Example,
    # 14=Key Design (was 12), 15=System Config (was 13), etc.
    # We need to update the title numbers
    renumber_map = {
        # index: (old_num_prefix, new_num_prefix)
        12: ("12", "12"),   # Gen Eval - already set as "12"
        13: ("13", "13"),   # Example - already set as "13"
        14: ("12", "14"),   # Key Design Decisions
        15: ("13", "15"),   # System Configuration
        16: ("14", "16"),   # Testing & Evaluation
        17: ("15", "17"),   # Recommendations
    }

    for idx, (old_n, new_n) in renumber_map.items():
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.startswith(f"{old_n}  "):
                        for run in para.runs:
                            if run.text.startswith(f"{old_n}  "):
                                run.text = run.text.replace(f"{old_n}  ", f"{new_n}  ", 1)
                                print(f"  Renumbered slide {idx+1}: {old_n} → {new_n}")
                                break
                        break

    # ── Save ──────────────────────────────────────────────────────────
    prs.save(OUTPUT_PATH)
    print(f"\n✅ Saved updated presentation to {OUTPUT_PATH}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
